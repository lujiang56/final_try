"""PPT 文件解析和文本提取 — 用于上传课件后的考点分析"""

import os
import re


ALLOWED_EXTENSIONS = {'.pptx', '.ppt'}
MAX_FILE_SIZE = 16 * 1024 * 1024  # 16 MB


def allowed_file(filename: str) -> bool:
    """检查文件扩展名是否在允许列表中"""
    _, ext = os.path.splitext(filename)
    return ext.lower() in ALLOWED_EXTENSIONS


def extract_text(filepath: str) -> tuple:
    """
    从 PPT/PPTX 文件中提取所有文本内容。

    Args:
        filepath: 上传文件在本地的路径

    Returns:
        (extracted_text, slide_count)
        extracted_text: 包含所有幻灯片文本的字符串，每页以分隔符标记
        slide_count: 幻灯片总页数
    """
    try:
        from pptx import Presentation
    except ImportError:
        return ("[错误] python-pptx 库未安装，请运行: pip install python-pptx", 0)

    try:
        prs = Presentation(filepath)
    except Exception as e:
        return (f"[错误] 无法打开 PPT 文件，文件可能已损坏或格式不支持: {e}", 0)

    all_text = []
    slide_count = len(prs.slides)

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []

        for shape in slide.shapes:
            # 提取文本框内容
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        slide_texts.append(para_text)

            # 提取表格内容
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        slide_texts.append(' | '.join(row_texts))

            # 提取组合形状内的文本 (递归)
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                _extract_group_text(shape, slide_texts)

        if slide_texts:
            all_text.append(f'\n--- 第{slide_num}页 ---\n' + '\n'.join(slide_texts))

    if not all_text:
        return ("[提示] 该 PPT 中未提取到文本内容，可能全部是图片。建议转换为带文本的格式。", slide_count)

    full_text = '\n'.join(all_text)
    return (full_text, slide_count)


def _extract_group_text(shape, slide_texts: list):
    """递归提取组合形状中的文本"""
    if hasattr(shape, 'shapes'):
        for child in shape.shapes:
            if child.has_text_frame:
                for para in child.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        slide_texts.append(para_text)
            if hasattr(child, 'has_table') and child.has_table:
                table = child.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        slide_texts.append(' | '.join(row_texts))
            if hasattr(child, 'shapes'):
                _extract_group_text(child, slide_texts)


def preprocess_text(raw_text: str) -> str:
    """
    清理和规范化 PPT 提取的文本，使其适配 analyzer 的分析流程。

    处理内容：
    1. 标记以数字+标点开头的行，作为"题目"分隔
    2. 压缩过多的空行
    3. 为纯知识点列表（无题目编号）添加人工题号

    Args:
        raw_text: PPT 提取的原始文本

    Returns:
        规范化后的文本，可直接传入 analyzer.analyze_materials()
    """
    if raw_text.startswith('[错误]') or raw_text.startswith('[提示]'):
        return raw_text

    lines = raw_text.split('\n')
    processed = []

    # 压缩连续空行
    prev_empty = False
    for line in lines:
        stripped = line.strip()
        if not stripped:
            if not prev_empty:
                processed.append('')
                prev_empty = True
        else:
            prev_empty = False
            # 如果行以数字+分隔符开头，保留作为题目标记
            processed.append(stripped)

    text = '\n'.join(processed)

    # 尝试识别存在知识点但无数字编号的情况
    # 如果文本中没有明显的题目编号 (如 "1."、"3、"、"一、")，给每个段落加数字编号
    has_numbering = bool(re.search(r'\n\s*\d+[\.\)、]', text)) or \
                     bool(re.search(r'\n\s*[一二三四五六七八九十]+[、．.]', text))

    if not has_numbering:
        # 将每段文本当作独立的知识点，添加序号
        paragraphs = [p.strip() for p in text.split('\n') if p.strip()]
        paragraphs = [p for p in paragraphs if not p.startswith('--- 第')]
        numbered = []
        for i, para in enumerate(paragraphs, 1):
            if len(para) > 5:  # 过滤过短的行
                numbered.append(f'{i}. {para}')
        if numbered:
            text = '\n'.join(numbered)

    return text


def extract_by_slide(filepath: str) -> list:
    """
    按幻灯片结构提取文本，保留每页的标题和内容。

    这是 PPT 分析的推荐入口——保留幻灯片结构信息，
    标题作为天然的章节标签。

    Args:
        filepath: PPT 文件路径

    Returns:
        list of dict: [
            {
                'slide_num': 1,
                'title': '第一章 数据表示与运算',
                'full_text': '全部文本...',
                'items': ['知识点1（10分）', '知识点2（15分）', ...],
                'has_table': bool,
            },
            ...
        ]
        如果解析失败，返回空列表
    """
    try:
        from pptx import Presentation
    except ImportError:
        return []

    try:
        prs = Presentation(filepath)
    except Exception:
        return []

    slides = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_data = {
            'slide_num': slide_num,
            'title': '',
            'full_text': '',
            'items': [],
            'has_table': False,
            'all_text_lines': [],
        }

        title_texts = []
        content_texts = []

        for shape in slide.shapes:
            is_title = getattr(shape, 'is_placeholder', False) and \
                       shape.placeholder_format.type == 1  # TITLE placeholder

            text_parts = []

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    para_text = para.text.strip()
                    if para_text:
                        text_parts.append(para_text)
                        if is_title:
                            title_texts.append(para_text)

            if shape.has_table:
                slide_data['has_table'] = True
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        cell_text = cell.text.strip()
                        if cell_text:
                            row_texts.append(cell_text)
                    if row_texts:
                        text_parts.append(' | '.join(row_texts))

            if shape.shape_type == 6:  # GROUP
                _extract_group_text(shape, text_parts)

            if text_parts:
                slide_data['all_text_lines'].extend(text_parts)
                if not is_title:
                    content_texts.extend(text_parts)

        # 判定标题：优先取 placeholder title，否则取第一个非空行
        if title_texts:
            slide_data['title'] = title_texts[0]
        elif slide_data['all_text_lines']:
            # 第一个行通常是标题（如果有 : 或 — 等标题特征）
            first = slide_data['all_text_lines'][0]
            if len(first) < 50 and not first.startswith('•') and not first.startswith('-'):
                slide_data['title'] = first
                slide_data['items'] = slide_data['all_text_lines'][1:]
            else:
                slide_data['items'] = slide_data['all_text_lines']

        if not slide_data['title']:
            slide_data['title'] = f'第{slide_num}页'
        if not slide_data['items']:
            slide_data['items'] = slide_data['all_text_lines']

        slide_data['full_text'] = '\n'.join(slide_data['all_text_lines'])

        # 过滤掉明显非知识点的行（如"考试情况总结"、"题型分布"等元信息）
        meta_keywords = ['题型分布', '考试情况', '建议策略', '高频考点（', '中等考点（', '低分考点（']
        slide_data['items'] = [
            item for item in slide_data['items']
            if not any(mk in item for mk in meta_keywords) or '分' in item
        ]
        # 但保留包含分值标记的项
        scored_meta = [
            item for item in slide_data['all_text_lines']
            if any(mk in item for mk in meta_keywords) and re.search(r'\d+\s*分', item)
        ]
        if scored_meta:
            slide_data['items'].extend(scored_meta)

        # 去重
        seen = set()
        unique_items = []
        for item in slide_data['items']:
            if item not in seen:
                seen.add(item)
                unique_items.append(item)
        slide_data['items'] = unique_items

        slides.append(slide_data)

    return slides
